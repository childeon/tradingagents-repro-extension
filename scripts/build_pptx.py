"""Rebuild FinalProjectPresentation.pptx with improved structure and content."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BASE    = os.path.dirname(os.path.abspath(__file__))
RESULTS = os.path.join(BASE, "phase1_results")
ASSETS  = os.path.join(BASE, "TradingAgents", "assets")

# ── colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1E, 0x3A, 0x5F)   # header bars, bold labels
BLUE   = RGBColor(0x2B, 0x6C, 0xB0)   # accents / section dividers
LTBLUE = RGBColor(0xBE, 0xD8, 0xF0)   # light blue tint (table header, lines)
BG     = RGBColor(0xF2, 0xF5, 0xFA)   # slide body background (clean cool white)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BODY   = RGBColor(0x2D, 0x37, 0x48)   # dark charcoal for body text (not pure black)
GRAY   = RGBColor(0x5A, 0x5F, 0x6E)   # secondary text
GREEN  = RGBColor(0x27, 0x6A, 0x3A)   # positive result
RED    = RGBColor(0xB8, 0x1C, 0x1C)   # negative result
AMBER  = RGBColor(0x92, 0x4E, 0x07)   # warning text

W = Inches(13.333)
H = Inches(7.5)

# Precomputed aspect ratios (width/height) from actual pixel dimensions:
# schema.png:              3470×1045  → 3.321
# equity_curves.png:       2179×980   → 2.223
# metrics_comparison.png:  2579×794   → 3.248


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def add_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    slide  = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    return slide


# ── primitive helpers ─────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=None, line=None, line_pt=0.75):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, size=14, bold=False, italic=False,
        color=BODY, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bullets(slide, items, l, t, w, h, size=13, color=BODY, sym="▸ "):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = sym + item
        r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def img(slide, path, l, t, w, ratio=None):
    """Add image preserving aspect ratio. Pass ratio=w/h to compute height."""
    if ratio:
        h = Inches(w.inches / ratio)
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        slide.shapes.add_picture(path, l, t, w)  # pptx auto-computes height


def header(slide, title, sub=None):
    rect(slide, 0, 0, W, Inches(1.05), fill=NAVY)
    # thin bright-blue accent line at very top
    rect(slide, 0, 0, W, Inches(0.055), fill=BLUE)
    txt(slide, title, Inches(0.4), Inches(0.07), Inches(12.5), Inches(0.6),
        size=27, bold=True, color=WHITE)
    if sub:
        txt(slide, sub, Inches(0.4), Inches(0.68), Inches(12.5), Inches(0.32),
            size=13, color=LTBLUE, italic=True)


def section_label(slide, text, l, t, w):
    """Narrow navy label bar above a content block."""
    rect(slide, l, t, w, Inches(0.30), fill=NAVY)
    txt(slide, text, l + Inches(0.12), t + Inches(0.04),
        w - Inches(0.12), Inches(0.26), size=10, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    slide = add_slide(prs)

    rect(slide, 0, 0, W, H, fill=NAVY)
    rect(slide, 0, 0, W, Inches(0.08), fill=BLUE)                 # top stripe
    rect(slide, 0, H - Inches(0.08), W, Inches(0.08), fill=BLUE)  # bottom stripe

    # subtle mid-slide rule
    rect(slide, Inches(1.2), Inches(3.28), Inches(10.9), Inches(0.035), fill=BLUE)

    txt(slide,
        "TradingAgents: Multi-Agent LLM Financial Trading Framework",
        Inches(1.0), Inches(1.3), Inches(11.3), Inches(1.7),
        size=37, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(slide, "Reproduction & Extension Study",
        Inches(1.0), Inches(3.42), Inches(11.3), Inches(0.55),
        size=22, italic=True, color=LTBLUE, align=PP_ALIGN.CENTER)

    txt(slide, "Machine Learning in Practice  ·  Columbia University  ·  Spring 2026",
        Inches(1.0), Inches(4.1), Inches(11.3), Inches(0.38),
        size=13, color=LTBLUE, align=PP_ALIGN.CENTER)

    txt(slide, "Yujia Zhang  ·  Ruiyan Li  ·  Yeyuxi Yi",
        Inches(1.0), Inches(4.55), Inches(11.3), Inches(0.38),
        size=13, color=LTBLUE, align=PP_ALIGN.CENTER)

    txt(slide, "Paper: arXiv 2412.20138  (Dec 2024)",
        Inches(1.0), Inches(5.25), Inches(11.3), Inches(0.32),
        size=11, color=RGBColor(0x7F, 0xAA, 0xD4), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Paper Overview
# ════════════════════════════════════════════════════════════════════════════
def slide_paper_overview(prs):
    slide = add_slide(prs)
    rect(slide, 0, 0, W, H, fill=BG)
    header(slide, "Paper Overview", "Why this paper · Task · Paper's contribution")

    Y0   = Inches(1.15)
    LPAD = Inches(0.4)
    COL1 = Inches(5.7)   # left column width (Why + Task)
    COL2 = Inches(6.8)   # right column width (Contribution)
    GAP  = Inches(0.3)
    X2   = LPAD + COL1 + GAP

    # ── left column ──────────────────────────────────────────────────────────
    y = Y0
    section_label(slide, "WHY THIS PAPER", LPAD, y, COL1)
    y += Inches(0.31)
    bullets(slide, [
        "Cutting-edge (Dec 2024): one of the first studies to replace specialist traders with LLMs",
        "Fuses heterogeneous signals: price, news, social media, fundamentals — all in one pipeline",
        "Structured bull/bear debate is a novel LLM design pattern worth stress-testing",
    ], LPAD, y, COL1, Inches(1.05), size=13)
    y += Inches(1.10)

    section_label(slide, "TASK", LPAD, y, COL1)
    y += Inches(0.31)
    txt(slide,
        "Fully automated daily stock-trading decisions for a single ticker driven by a "
        "multi-agent LLM pipeline — no hand-coded signals, no human intervention.",
        LPAD, y, COL1, Inches(0.68), size=13, color=GRAY)

    # ── right column ─────────────────────────────────────────────────────────
    y = Y0
    section_label(slide, "PAPER'S CONTRIBUTION", X2, y, COL2)
    y += Inches(0.31)
    bullets(slide, [
        "4 specialised teams: Analyst → Researcher → Trader → Risk Manager",
        "Bull vs. Bear debate forces the model to weigh opposing evidence before deciding",
        "Risk Management layer moderates position sizing before order execution",
        "Beats MACD, SMA, KDJ+RSI, and Buy-and-Hold on cumulative return & Sharpe ratio",
    ], X2, y, COL2, Inches(1.3), size=13)

    # ── schema diagram (full-width, bottom) ──────────────────────────────────
    # schema.png  3470×1045  ratio 3.321
    # width = 10.5" → height = 10.5/3.321 = 3.16"  → bottom at y=4.2+3.16=7.36 ✓
    schema = os.path.join(ASSETS, "schema.png")
    img(slide, schema, Inches(1.42), Inches(4.2), Inches(10.5), ratio=3.321)

    # light rule between text and schema
    rect(slide, LPAD, Inches(4.12), W - LPAD * 2, Inches(0.02), fill=LTBLUE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Reproduction: Scope & Challenges
# ════════════════════════════════════════════════════════════════════════════
def slide_reproduction(prs):
    slide = add_slide(prs)
    rect(slide, 0, 0, W, H, fill=BG)
    header(slide, "Reproduction", "Scope reduction · Challenges · Solutions")

    LPAD = Inches(0.4)
    FULL = W - LPAD * 2

    # ── scope banner ──────────────────────────────────────────────────────────
    rect(slide, LPAD, Inches(1.15), FULL, Inches(1.02),
         fill=WHITE, line=LTBLUE, line_pt=0.75)
    txt(slide, "Scope Reduction",
        LPAD + Inches(0.18), Inches(1.2), FULL, Inches(0.3),
        size=12, bold=True, color=NAVY)
    txt(slide,
        "Original paper:  3 tickers × 3 months   →   Our replication:  AMZN only, "
        "January 2024  (21 trading days)\n"
        "Reason: full replication required prohibitive LLM API cost and runtime",
        LPAD + Inches(0.18), Inches(1.51), FULL - Inches(0.3), Inches(0.58),
        size=12.5, color=GRAY)

    # ── challenges ────────────────────────────────────────────────────────────
    challenges = [
        (NAVY,
         "Challenge 1 — Stale Fundamentals Data",
         "The Fundamentals Analyst pulled the most-recent SEC filing available at run time. "
         "Running in 2026, the API returned 2026 filings — not 2024 as needed for the backtest.",
         "Fix: Manually downloaded the Q3 2023 10-Q and cached it for the entire backtest window."),
        (BLUE,
         "Challenge 2 — News API Ordering Bug",
         "yfinance returns news chronologically-latest first, then filters by date. Running in 2026 "
         "with a 2024 date filter returned zero results — the filter ran after the API cutoff, not before.",
         "Fix: Replaced live API calls with pre-filtered static GDELT + Reddit CSV datasets for Jan 2024."),
        (RGBColor(0x44, 0x72, 0xC4),
         "Challenge 3 — Latency",
         "Sequential LLM calls per trading day (analyst → researcher → trader → risk) "
         "were expensive in both time and API cost.",
         "Fix: Cached the fundamentals report across the full month (it's a single quarter's filing)."),
    ]

    y = Inches(2.28)
    for bar_color, title, problem, fix in challenges:
        # title bar
        rect(slide, LPAD, y, FULL, Inches(0.29), fill=bar_color)
        txt(slide, title, LPAD + Inches(0.12), y + Inches(0.04),
            FULL, Inches(0.25), size=11.5, bold=True, color=WHITE)
        y += Inches(0.30)
        # body (two-line: problem / fix)
        rect(slide, LPAD, y, FULL, Inches(0.85), fill=WHITE, line=LTBLUE, line_pt=0.5)
        txt(slide, f"Problem: {problem}",
            LPAD + Inches(0.18), y + Inches(0.06), FULL - Inches(0.25), Inches(0.38),
            size=11.5, color=GRAY)
        txt(slide, f"Fix: {fix}",
            LPAD + Inches(0.18), y + Inches(0.44), FULL - Inches(0.25), Inches(0.36),
            size=11.5, bold=True, color=BODY)
        y += Inches(0.90)

    # limitation footnote
    rect(slide, LPAD, Inches(6.92), FULL, Inches(0.44),
         fill=RGBColor(0xFF, 0xF3, 0xCD), line=RGBColor(0xE0, 0xA8, 0x00), line_pt=0.75)
    txt(slide,
        "⚠  Limitation: static datasets make the system non-production-ready. "
        "The paper's live-API design is the intended architecture; ours is a faithful but frozen snapshot.",
        LPAD + Inches(0.15), Inches(6.96), FULL - Inches(0.2), Inches(0.38),
        size=11, italic=True, color=AMBER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Replication Results
# ════════════════════════════════════════════════════════════════════════════
def slide_results(prs):
    slide = add_slide(prs)
    rect(slide, 0, 0, W, H, fill=BG)
    header(slide, "Replication Results", "AMZN · January 2024 · 21 trading days")

    LPAD = Inches(0.4)

    # ── metrics table (left) ──────────────────────────────────────────────────
    headers_  = ["Strategy", "Cum. Return", "Ann. Return", "Sharpe", "Max DD"]
    col_ws    = [Inches(1.7), Inches(0.95), Inches(0.95), Inches(0.82), Inches(0.88)]
    rows = [
        ("TradingAgents", "+4.41%", "+67.8%", "3.90", "−3.57%", True),
        ("Buy & Hold",    "+3.51%", "+51.4%", "2.40", "−3.76%", False),
        ("KDJ + RSI",     "+3.26%", "+46.9%", "3.54", "−1.34%", False),
        ("MACD",          "+0.96%", "+12.1%", "0.82", "−3.76%", False),
        ("ZMR / SMA",     "0.00%",  "0.00%",  "0.00", "0.00%",  False),
    ]
    ROW_H = Inches(0.38)
    tt = Inches(1.18)

    # header row
    x = LPAD
    for hdr, cw in zip(headers_, col_ws):
        rect(slide, x, tt, cw, Inches(0.33), fill=NAVY)
        txt(slide, hdr, x + Inches(0.05), tt + Inches(0.04),
            cw - Inches(0.05), Inches(0.28), size=10, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        x += cw

    for ri, (strat, cr, ar, sh, md, hi) in enumerate(rows):
        y_row = tt + Inches(0.33) + ri * ROW_H
        bg    = LTBLUE if hi else (WHITE if ri % 2 == 0 else RGBColor(0xF0, 0xF4, 0xF9))
        x     = LPAD
        for ci, (val, cw) in enumerate(zip([strat, cr, ar, sh, md], col_ws)):
            rect(slide, x, y_row, cw, ROW_H, fill=bg,
                 line=RGBColor(0xD4, 0xDB, 0xE8), line_pt=0.4)
            fcol = NAVY if hi else (BODY if ci == 0 else GRAY)
            txt(slide, val, x + Inches(0.05), y_row + Inches(0.06),
                cw - Inches(0.05), ROW_H - Inches(0.06),
                size=11, bold=(hi and ci == 0), color=fcol,
                align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
            x += cw

    table_bottom = tt + Inches(0.33) + len(rows) * ROW_H   # ≈ 3.71"

    # takeaway
    rect(slide, LPAD, table_bottom + Inches(0.1), Inches(5.3), Inches(0.46),
         fill=RGBColor(0xDC, 0xF0, 0xE3), line=GREEN, line_pt=0.75)
    txt(slide, "✓  TradingAgents beats every benchmark on return and Sharpe ratio",
        LPAD + Inches(0.12), table_bottom + Inches(0.17), Inches(5.1), Inches(0.32),
        size=11, bold=True, color=GREEN)

    # ── equity curves (right of table) ────────────────────────────────────────
    # equity_curves.png  2179×980  ratio 2.223
    # width = 7.35" → height = 7.35/2.223 = 3.31"
    eq_x = LPAD + Inches(5.5)
    eq_w = Inches(7.35)
    img(slide, os.path.join(RESULTS, "equity_curves.png"),
        eq_x, Inches(1.18), eq_w, ratio=2.223)

    # ── metrics comparison (bottom, centred) ──────────────────────────────────
    # metrics_comparison.png  2579×794  ratio 3.248
    # width = 9.3" → height = 9.3/3.248 = 2.86"  → bottom at 4.6+2.86=7.46 ✓
    mc_w = Inches(9.3)
    mc_x = (W - mc_w) / 2
    img(slide, os.path.join(RESULTS, "metrics_comparison.png"),
        mc_x, Inches(4.58), mc_w, ratio=3.248)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Extension: Ablation Study
# ════════════════════════════════════════════════════════════════════════════
def slide_extension(prs):
    slide = add_slide(prs)
    rect(slide, 0, 0, W, H, fill=BG)
    header(slide, "Extension: Ablation Study",
           "What happens when we remove individual agent layers?")

    LPAD = Inches(0.4)
    FULL = W - LPAD * 2

    # ── summary strip ─────────────────────────────────────────────────────────
    rect(slide, LPAD, Inches(1.12), FULL, Inches(0.35), fill=WHITE, line=LTBLUE, line_pt=0.5)
    configs = [
        ("Full Architecture  +4.41%  Sharpe 3.90",   BLUE),
        ("No Research Debate  +6.31%  Sharpe 5.73",  GREEN),
        ("No Risk Layer  −1.33%  Sharpe −0.77",       RED),
        ("Buy & Hold  +3.51%  Sharpe 2.40",           GRAY),
    ]
    cx = LPAD + Inches(0.15)
    for label, col in configs:
        txt(slide, label, cx, Inches(1.17), Inches(3.1), Inches(0.26),
            size=10, bold=False, color=col)
        cx += Inches(3.22)

    # ── Experiment 1 ─────────────────────────────────────────────────────────
    EX_W   = W - LPAD * 2
    NUM_W  = Inches(2.2)
    TXT_W  = EX_W - NUM_W - Inches(0.2)
    TXT_X  = LPAD + NUM_W + Inches(0.2)
    EX1_Y  = Inches(1.57)

    # green header
    rect(slide, LPAD, EX1_Y, EX_W, Inches(0.32), fill=GREEN)
    txt(slide, "Experiment 1 — Remove Bull/Bear Research Debate",
        LPAD + Inches(0.15), EX1_Y + Inches(0.05), EX_W, Inches(0.27),
        size=12.5, bold=True, color=WHITE)

    EX1_Y += Inches(0.34)
    # result callout
    rect(slide, LPAD, EX1_Y, NUM_W, Inches(1.2),
         fill=RGBColor(0xE6, 0xF4, 0xEC), line=GREEN, line_pt=1.0)
    txt(slide, "+6.31%", LPAD, EX1_Y + Inches(0.1), NUM_W, Inches(0.58),
        size=30, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(slide, "Sharpe  5.73", LPAD, EX1_Y + Inches(0.67), NUM_W, Inches(0.28),
        size=13, color=GREEN, align=PP_ALIGN.CENTER)
    txt(slide, "▲ BETTER", LPAD, EX1_Y + Inches(0.93), NUM_W, Inches(0.24),
        size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    txt(slide,
        "The bull/bear debate introduced noise and indecision into an otherwise trending market. "
        "In January 2024, AMZN was in a clear uptrend; the bearish debater consistently "
        "argued against well-founded bullish signals, forcing the trader to hedge unnecessarily. "
        "Removing the debate layer lets the Researcher pass confident, unfiltered "
        "recommendations directly — yielding stronger performance when market direction is clear.",
        TXT_X, EX1_Y + Inches(0.08), TXT_W, Inches(1.1),
        size=12.5, color=GRAY)

    # ── Experiment 2 ─────────────────────────────────────────────────────────
    EX2_Y = EX1_Y + Inches(1.35)

    rect(slide, LPAD, EX2_Y, EX_W, Inches(0.32), fill=RED)
    txt(slide, "Experiment 2 — Remove Risk Management Layer",
        LPAD + Inches(0.15), EX2_Y + Inches(0.05), EX_W, Inches(0.27),
        size=12.5, bold=True, color=WHITE)

    EX2_Y += Inches(0.34)
    rect(slide, LPAD, EX2_Y, NUM_W, Inches(1.2),
         fill=RGBColor(0xFC, 0xE8, 0xE8), line=RED, line_pt=1.0)
    txt(slide, "−1.33%", LPAD, EX2_Y + Inches(0.1), NUM_W, Inches(0.58),
        size=30, bold=True, color=RED, align=PP_ALIGN.CENTER)
    txt(slide, "Sharpe −0.77", LPAD, EX2_Y + Inches(0.67), NUM_W, Inches(0.28),
        size=13, color=RED, align=PP_ALIGN.CENTER)
    txt(slide, "▼ MUCH WORSE", LPAD, EX2_Y + Inches(0.93), NUM_W, Inches(0.24),
        size=10, bold=True, color=RED, align=PP_ALIGN.CENTER)

    txt(slide,
        "The Risk Management team is the critical guardrail of the pipeline. "
        "Without it, the Trader executes aggressive, unchecked positions based on raw "
        "Researcher output. Even small over-confident bets compound: max drawdown worsens "
        "from −3.57% to −4.81% and cumulative return goes negative. "
        "This confirms that risk management is not redundant overhead — it is load-bearing.",
        TXT_X, EX2_Y + Inches(0.08), TXT_W, Inches(1.1),
        size=12.5, color=GRAY)

    # ── key insight footer ────────────────────────────────────────────────────
    FT_Y = EX2_Y + Inches(1.34)
    rect(slide, LPAD, FT_Y, EX_W, Inches(0.5),
         fill=RGBColor(0xE8, 0xED, 0xF7), line=NAVY, line_pt=0.75)
    txt(slide,
        "Key Insight: Debate adds value only when market direction is uncertain. "
        "Risk management is always essential. Selectively disabling debate on trending markets "
        "may be a simple, practical way to improve framework performance.",
        LPAD + Inches(0.15), FT_Y + Inches(0.08), EX_W - Inches(0.3), Inches(0.4),
        size=11.5, bold=True, color=NAVY)


# ════════════════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════════════════
prs = new_prs()
slide_title(prs)
slide_paper_overview(prs)
slide_reproduction(prs)
slide_results(prs)
slide_extension(prs)

out = os.path.join(BASE, "FinalProjectPresentation_v2.pptx")
prs.save(out)
print("Saved:", out)
